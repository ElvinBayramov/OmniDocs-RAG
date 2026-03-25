"""
Markdown RAG MCP Server v3.3 (Windows-compatible)

Features:
- Multi-format support (md, txt, rst, html, code files)
- Multilingual embeddings (BAAI/bge-m3, GPU-accelerated if available)
- Hybrid Search: ChromaDB vector + BM25 keyword scoring
- Cross-Encoder Reranking (multilingual, configurable)
- Incremental indexing (only re-indexes changed files)
- Multi-collection support
- BM25 persistence across server restarts
- Metadata filters (category, filename)
- Admin tools (list collections, remove sources)
"""

import os
import re
import json
import glob
import pickle
import warnings
import logging
import hashlib
from pathlib import Path
from typing import Optional

os.environ["HF_HUB_DISABLE_PROGRESS_BARS"] = "1"
os.environ["TOKENIZERS_PARALLELISM"] = "false"
os.environ["TRANSFORMERS_VERBOSITY"] = "error"
warnings.filterwarnings("ignore")
logging.basicConfig(level=logging.ERROR)
logging.getLogger("sentence_transformers").setLevel(logging.ERROR)
logging.getLogger("transformers").setLevel(logging.ERROR)

from fastmcp import FastMCP
import chromadb
from chromadb.utils import embedding_functions

# ──────────────────────────────────────────────
# Config
# ──────────────────────────────────────────────
_SERVER_DIR = Path(__file__).parent
DOCS_PATH = os.getenv(
    "RAG_DOCS_PATH",
    str(_SERVER_DIR.parent)
)
DB_PATH = os.getenv(
    "RAG_DB_PATH",
    str(_SERVER_DIR / "chroma_db")
)
COLLECTION_NAME = "docs_v4"
DEFAULT_COLLECTION = COLLECTION_NAME
EMBED_MODEL = os.getenv("RAG_EMBED_MODEL", "BAAI/bge-m3")
RERANK_MODEL = os.getenv("RAG_RERANK_MODEL", "BAAI/bge-reranker-v2-m3")
MAX_CROSS_ENCODER_CHARS = 8000   # bge-reranker-v2-m3 supports up to 8192 tokens
MAX_RESULT_CHARS = 1500          # output truncation limit per result
_HASH_CACHE_PATH = _SERVER_DIR / "data" / "file_hashes.json"

# ──────────────────────────────────────────────
# Device detection (GPU optional — no torch required)
# ──────────────────────────────────────────────
def _detect_device() -> str:
    """Detect compute device. Set RAG_DEVICE=cuda to force GPU."""
    requested = os.getenv("RAG_DEVICE", "auto").lower()
    if requested == "cpu":
        return "cpu"
    try:
        import torch
        if requested == "cuda" or (requested == "auto" and torch.cuda.is_available()):
            return "cuda"
    except ImportError:
        pass
    return "cpu"

DEVICE = _detect_device()

# ──────────────────────────────────────────────
# ChromaDB + Models setup
# ──────────────────────────────────────────────
client = chromadb.PersistentClient(path=DB_PATH)

# GPU-aware embedding function with normalization (required for bge-m3)
try:
    from sentence_transformers import SentenceTransformer
    _embed_model = SentenceTransformer(EMBED_MODEL, device=DEVICE)

    class _EmbeddingFunction:
        """ChromaDB-compatible embedding function with normalization + optional GPU."""
        is_legacy = True  # Tell ChromaDB 0.6.0 to treat this as old-style callable
        
        def name(self) -> str:
            return EMBED_MODEL

        def __call__(self, input: list[str]) -> list[list[float]]:
            return _embed_model.encode(
                input,
                batch_size=32,
                show_progress_bar=False,
                convert_to_numpy=True,
                normalize_embeddings=True,  # required for bge-m3 cosine similarity
            ).tolist()
            
        def embed_documents(self, input: list[str]) -> list[list[float]]:
            return self(input)
            
        def embed_query(self, input: list[str]) -> list[list[float]]:
            return self(input)


    embed_fn = _EmbeddingFunction()
except Exception:
    # Fallback: ChromaDB built-in (no normalization, CPU only)
    embed_fn = embedding_functions.SentenceTransformerEmbeddingFunction(
        model_name=EMBED_MODEL
    )
    embed_fn.is_legacy = True  # safety patch

collection = client.get_or_create_collection(
    name=COLLECTION_NAME,
    embedding_function=embed_fn,
    metadata={"hnsw:space": "cosine"}
)

# Lazy load cross-encoder (downloads ~1.1GB on first use)
_cross_encoder = None

def _get_cross_encoder():
    global _cross_encoder
    if _cross_encoder is None:
        from sentence_transformers import CrossEncoder
        _cross_encoder = CrossEncoder(RERANK_MODEL)
    return _cross_encoder

# BM25 index (built after indexing, persisted to disk)
_bm25_index = None
_bm25_corpus = None   # list of (chunk_id, text)
_BM25_CACHE_PATH = _SERVER_DIR / "data" / "bm25_cache.pkl"

def _tokenize(text: str) -> list[str]:
    """Simple whitespace + lowercase tokenizer for BM25."""
    text = re.sub(r"[^\w\s]", " ", text.lower())
    return [w for w in text.split() if len(w) > 2]

def _build_bm25(ids: list[str], texts: list[str]):
    """Build BM25 index and persist to disk for restart survival."""
    global _bm25_index, _bm25_corpus
    if not ids:
        _bm25_index = None
        _bm25_corpus = None
        return
    from rank_bm25 import BM25Okapi
    tokenized = [_tokenize(t) for t in texts]
    _bm25_index = BM25Okapi(tokenized)
    _bm25_corpus = list(zip(ids, texts))
    # Persist to disk so BM25 survives server restarts
    try:
        _BM25_CACHE_PATH.parent.mkdir(parents=True, exist_ok=True)
        with open(_BM25_CACHE_PATH, "wb") as f:
            pickle.dump({"corpus": _bm25_corpus, "tokenized": tokenized}, f)
    except Exception:
        pass  # non-critical: BM25 still works in-memory

_bm25_loaded_from = None  # tracks how BM25 was initialized

def _load_bm25_on_startup():
    """Restore BM25 index from disk cache if available and still valid."""
    global _bm25_index, _bm25_corpus, _bm25_loaded_from
    if not _BM25_CACHE_PATH.exists():
        return
    try:
        from rank_bm25 import BM25Okapi
        with open(_BM25_CACHE_PATH, "rb") as f:
            data = pickle.load(f)
        # Validate: cache must match current ChromaDB collection size
        chroma_count = collection.count()
        cache_count = len(data["corpus"])
        if chroma_count > 0 and cache_count != chroma_count:
            return  # stale cache — skip, will rebuild on next index_documents()
        _bm25_index = BM25Okapi(data["tokenized"])
        _bm25_corpus = data["corpus"]
        _bm25_loaded_from = "cache"
    except Exception:
        pass  # corrupt cache — will rebuild on next index_documents()

_load_bm25_on_startup()

mcp = FastMCP("Markdown RAG")


# ──────────────────────────────────────────────
# Heading-Aware Chunking with Overlap + Code-Aware
# ──────────────────────────────────────────────
def _split_into_sentences(text: str) -> list[str]:
    """Split text into sentences for overlap."""
    pieces = re.split(r"(?<=[.!?。])\s+|\n{2,}", text)
    return [s.strip() for s in pieces if s.strip()]

def _extract_sections(text: str, filepath: str) -> list[dict]:
    """Split markdown by headings (## and ###) into semantic chunks
    with 2-sentence overlap between consecutive chunks."""

    # Strip YAML frontmatter (\A = absolute start of string, safer than ^)
    text = re.sub(r"\A---.*?---\s*", "", text, flags=re.DOTALL)

    parts = re.split(r"(^#{1,3}\s+.+$)", text, flags=re.MULTILINE)

    raw_sections = []
    current_heading = "Introduction"
    current_h1 = Path(filepath).stem

    for part in parts:
        part = part.strip()
        if not part:
            continue

        heading_match = re.match(r"^(#{1,3})\s+(.+)$", part)
        if heading_match:
            level = len(heading_match.group(1))
            title = heading_match.group(2).strip()
            if level == 1:
                current_h1 = title
            current_heading = title
            continue

        clean = re.sub(r"[|\-\s#*`>]", "", part)
        if len(clean) < 30:
            continue

        # Split long sections into sub-chunks (~600 words to utilize bge-m3's 8192-token capacity)
        words = part.split()
        if len(words) > 700:
            sub_chunks = []
            for i in range(0, len(words), 600):
                sub_chunk = " ".join(words[i:i + 600])
                sub_chunks.append(sub_chunk)
        else:
            sub_chunks = [part]

        for idx, chunk_text in enumerate(sub_chunks):
            raw_sections.append({
                "text": chunk_text,
                "heading": current_heading,
                "parent_heading": current_h1,
                "source": filepath,
                "filename": Path(filepath).name,
                "sub_index": idx,
                "word_count": len(chunk_text.split())
            })

    # Add 2-sentence overlap between consecutive chunks from same file
    for i in range(1, len(raw_sections)):
        if raw_sections[i]["source"] == raw_sections[i-1]["source"]:
            prev_sentences = _split_into_sentences(raw_sections[i-1]["text"])
            overlap = " ".join(prev_sentences[-2:]) if len(prev_sentences) >= 2 else ""
            if overlap:
                raw_sections[i]["text"] = f"[...] {overlap}\n\n{raw_sections[i]['text']}"
                raw_sections[i]["word_count"] = len(raw_sections[i]["text"].split())

    return raw_sections


def _chunk_python_code(text: str, filepath: str) -> list[dict]:
    """Split Python source by top-level classes and functions via AST."""
    import ast
    chunks = []
    try:
        tree = ast.parse(text)
        lines = text.split("\n")
        for node in ast.iter_child_nodes(tree):
            if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
                start = node.lineno - 1
                end = node.end_lineno
                chunk_text = "\n".join(lines[start:end])
                node_type = type(node).__name__.replace("Def", "").replace("Async", "async ")
                chunks.append({
                    "text": f"```python\n{chunk_text}\n```",
                    "heading": f"{node_type}: {node.name}",
                    "parent_heading": Path(filepath).stem,
                    "source": filepath,
                    "filename": Path(filepath).name,
                    "sub_index": 0,
                    "word_count": len(chunk_text.split()),
                })
    except SyntaxError:
        return []  # fallback to standard chunking
    return chunks


def _chunk_js_code(text: str, filepath: str) -> list[dict]:
    """Split JS/TS source by function/class declarations via regex."""
    patterns = [
        r"((?:export\s+)?(?:async\s+)?function\s+\w+[^}]*\{(?:[^{}]|\{[^{}]*\})*\})",
        r"((?:export\s+)?class\s+\w+[^}]*\{(?:[^{}]|\{(?:[^{}]|\{[^{}]*\})*\})*\})",
        r"((?:export\s+)?(?:const|let|var)\s+\w+\s*=\s*(?:async\s+)?\([^)]*\)\s*=>\s*\{(?:[^{}]|\{[^{}]*\})*\})",
    ]
    chunks = []
    seen_ranges: list[tuple[int, int]] = []

    for pattern in patterns:
        for match in re.finditer(pattern, text, re.DOTALL):
            start, end = match.start(), match.end()
            # Skip if overlaps with already-found chunk
            if any(s <= start < e for s, e in seen_ranges):
                continue
            seen_ranges.append((start, end))
            chunk_text = match.group(1)
            # Extract name
            name_match = re.search(r"(?:function|class|const|let|var)\s+(\w+)", chunk_text)
            name = name_match.group(1) if name_match else "anonymous"
            node_type = "class" if "class " in chunk_text[:20] else "function"
            lang = Path(filepath).suffix.lstrip(".")
            chunks.append({
                "text": f"```{lang}\n{chunk_text}\n```",
                "heading": f"{node_type}: {name}",
                "parent_heading": Path(filepath).stem,
                "source": filepath,
                "filename": Path(filepath).name,
                "sub_index": 0,
                "word_count": len(chunk_text.split()),
            })

    return chunks


def _extract_sections_smart(text: str, filepath: str) -> list[dict]:
    """
    Smart chunking by file type:
    - .py     → by classes/functions via AST
    - .js/.ts → by functions via regex
    - others  → by headings (standard behavior)
    """
    suffix = Path(filepath).suffix.lower()

    if suffix == ".py":
        result = _chunk_python_code(text, filepath)
        if result:
            return result
    elif suffix in (".js", ".ts", ".jsx", ".tsx"):
        result = _chunk_js_code(text, filepath)
        if result:
            return result

    return _extract_sections(text, filepath)


def _categorize_file(filepath: str, content: str) -> str:
    """
    Auto-categorize file — three-level priority:
    1. YAML Frontmatter `category:` key (explicit user override)
    2. First H1 heading `# Title` in the document (dynamic, zero-effort)
    3. Filename stem as last resort
    """
    # Priority 1: YAML Frontmatter
    frontmatter_match = re.match(r"^---\s*\n(.*?)\n---", content, flags=re.DOTALL)
    if frontmatter_match:
        for line in frontmatter_match.group(1).split('\n'):
            if line.strip().lower().startswith('category:'):
                return line.split(':', 1)[1].strip().lower()

    # Priority 2: First H1 heading — "# My Document Title" → "my document title"
    h1_match = re.search(r"^#\s+(.+)$", content, flags=re.MULTILINE)
    if h1_match:
        raw_title = h1_match.group(1).strip()
        # Normalize: lowercase, strip markdown emphasis, keep letters/numbers/spaces
        category = re.sub(r"[*_`]", "", raw_title).lower()
        category = re.sub(r"[^\w\s-]", "", category).strip()
        if category:
            return category

    # Priority 3: Filename stem
    return Path(filepath).stem.lower().replace("_", " ").replace("-", " ")


def _format_result(doc: str, meta: dict, score: float, rank: int) -> str:
    """Format a single search result cleanly."""
    heading = meta.get("heading", "")
    parent = meta.get("parent_heading", "")
    filename = meta.get("filename", "unknown")
    category = meta.get("category", "")

    breadcrumb = filename
    if parent and parent != Path(filename).stem:
        breadcrumb += f" > {parent}"
    if heading and heading != parent:
        breadcrumb += f" > {heading}"

    content = doc.strip()
    # Remove overlap prefix markers
    content = re.sub(r"^\[\.\.\.\]\s*", "", content)
    content = re.sub(r"\n{3,}", "\n\n", content)

    if len(content) > MAX_RESULT_CHARS:
        # Don't truncate inside ```code blocks```
        code_block_end = content.find("```", MAX_RESULT_CHARS - 200)
        if 0 < code_block_end < MAX_RESULT_CHARS + 500:
            content = content[:code_block_end + 3]
        else:
            cut = content[:MAX_RESULT_CHARS].rfind(". ")
            if cut > MAX_RESULT_CHARS // 2:
                content = content[:cut + 1]
            else:
                content = content[:MAX_RESULT_CHARS] + "..."

    return (
        f"### [{rank}] {breadcrumb}\n"
        f"**Relevance:** {score:.0%} | **Category:** {category}\n\n"
        f"{content}\n"
    )


# ──────────────────────────────────────────────
# Multi-Format File Reader
# ──────────────────────────────────────────────
SUPPORTED_EXTENSIONS = {
    # Native text formats
    ".md", ".txt", ".rst", ".text", ".log",
    # Code files (wrapped in markdown for chunking)
    ".py", ".js", ".ts", ".jsx", ".tsx", ".css", ".scss",
    ".java", ".go", ".rs", ".c", ".cpp", ".h", ".hpp",
    ".rb", ".php", ".swift", ".kt", ".lua", ".sh", ".bash",
    # Data/config formats
    ".json", ".yaml", ".yml", ".toml", ".xml", ".csv", ".ini", ".cfg",
    # Web formats
    ".html", ".htm",
    # Binary document formats (require optional packages)
    ".pdf", ".docx", ".xlsx", ".pptx", ".ipynb",
}

_CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".css", ".scss",
    ".java", ".go", ".rs", ".c", ".cpp", ".h", ".hpp",
    ".rb", ".php", ".swift", ".kt", ".lua", ".sh", ".bash",
    ".json", ".yaml", ".yml", ".toml", ".xml", ".ini", ".cfg",
}

def _read_file_to_text(filepath: str) -> Optional[str]:
    """Read a file and convert to markdown-like text for indexing."""
    ext = Path(filepath).suffix.lower()
    try:
        raw = Path(filepath).read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return None

    if not raw.strip():
        return None

    # Markdown / plain text — use as-is
    if ext in (".md", ".txt", ".text", ".rst", ".log"):
        return raw

    # HTML — strip tags to plain text
    if ext in (".html", ".htm"):
        try:
            from html.parser import HTMLParser
            class _StripHTML(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.parts = []
                    self._skip = False
                def handle_starttag(self, tag, attrs):
                    if tag in ("script", "style"):
                        self._skip = True
                def handle_endtag(self, tag):
                    if tag in ("script", "style"):
                        self._skip = False
                def handle_data(self, data):
                    if not self._skip:
                        self.parts.append(data)
            parser = _StripHTML()
            parser.feed(raw)
            return "\n".join(parser.parts)
        except Exception:
            return re.sub(r"<[^>]+>", "", raw)  # fallback: naive strip

    # Code files — wrap in markdown code block with filename heading
    if ext in _CODE_EXTENSIONS:
        lang = ext.lstrip(".")
        name = Path(filepath).name
        return f"# {name}\n\n```{lang}\n{raw}\n```"

    # ── Binary document formats (optional packages) ──

    # PDF via pypdf
    if ext == ".pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(filepath)
            pages = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"## Page {i + 1}\n\n{text}")
            return "\n\n".join(pages) if pages else None
        except ImportError:
            return None  # pypdf not installed
        except Exception:
            return None

    # Word (.docx) via python-docx
    if ext == ".docx":
        try:
            import docx
            doc = docx.Document(filepath)
            lines = []
            for para in doc.paragraphs:
                style = para.style.name
                if style.startswith("Heading 1"):
                    lines.append(f"# {para.text}")
                elif style.startswith("Heading 2"):
                    lines.append(f"## {para.text}")
                elif style.startswith("Heading 3"):
                    lines.append(f"### {para.text}")
                elif para.text.strip():
                    lines.append(para.text)
            # Extract tables as markdown
            for table in doc.tables:
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if rows:
                    header = "| " + " | ".join(rows[0]) + " |"
                    sep = "|" + "|".join(["---"] * len(rows[0])) + "|"
                    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
                    lines.append(f"{header}\n{sep}\n{body}")
            return "\n\n".join(lines) if lines else None
        except ImportError:
            return None  # python-docx not installed
        except Exception:
            return None

    # Excel (.xlsx) via openpyxl
    if ext == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, data_only=True)
            sections = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        rows.append([str(c) if c is not None else "" for c in row])
                if rows:
                    header = "| " + " | ".join(rows[0]) + " |"
                    sep = "|" + "|".join(["---"] * len(rows[0])) + "|"
                    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:50])
                    sections.append(f"## Sheet: {sheet_name}\n\n{header}\n{sep}\n{body}")
            return "\n\n".join(sections) if sections else None
        except ImportError:
            return None  # openpyxl not installed
        except Exception:
            return None

    # PowerPoint (.pptx) via python-pptx
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text.strip() for shape in slide.shapes
                         if hasattr(shape, "text") and shape.text.strip()]
                if texts:
                    slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
            return "\n\n".join(slides) if slides else None
        except ImportError:
            return None  # python-pptx not installed
        except Exception:
            return None

    # Jupyter Notebook (.ipynb)
    if ext == ".ipynb":
        try:
            nb = json.loads(Path(filepath).read_text(encoding="utf-8"))
            sections = []
            for cell in nb.get("cells", []):
                source = "".join(cell.get("source", []))
                if not source.strip():
                    continue
                if cell["cell_type"] == "markdown":
                    sections.append(source)
                elif cell["cell_type"] == "code":
                    sections.append(f"```python\n{source}\n```")
            return "\n\n".join(sections) if sections else None
        except Exception:
            return None

    return raw


# ──────────────────────────────────────────────
# Multi-Collection Support
# ──────────────────────────────────────────────
def _get_collection(name: str = DEFAULT_COLLECTION):
    """Get or create a named ChromaDB collection."""
    return client.get_or_create_collection(
        name=name,
        embedding_function=embed_fn,
        metadata={"hnsw:space": "cosine"}
    )


# ──────────────────────────────────────────────
# MCP Tools
# ──────────────────────────────────────────────
# ──────────────────────────────────────────────
# File hash cache for incremental indexing
# ──────────────────────────────────────────────
def _load_hash_cache() -> dict:
    if _HASH_CACHE_PATH.exists():
        try:
            return json.loads(_HASH_CACHE_PATH.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {}

def _save_hash_cache(cache: dict):
    _HASH_CACHE_PATH.parent.mkdir(parents=True, exist_ok=True)
    _HASH_CACHE_PATH.write_text(json.dumps(cache, indent=2), encoding="utf-8")

def _get_file_hash(filepath: str) -> str:
    return hashlib.md5(Path(filepath).read_bytes()).hexdigest()


@mcp.tool()
def index_documents(
    docs_path: str = DOCS_PATH,
    collection: str = DEFAULT_COLLECTION,
    force_reindex: bool = False
) -> str:
    """
    Index documents in the given directory into ChromaDB.
    Supports: .md, .txt, .rst, .html, .py, .js, .ts, and more.
    Uses incremental indexing — only re-indexes files that have changed.

    Args:
        docs_path: Path to scan for files (default: parent directory)
        collection: Collection name to index into (default: docs_v4)
        force_reindex: If True, ignore hash cache and re-index everything
    """
    # Gather all supported files
    all_files = []
    docs_p = Path(docs_path).resolve()
    for ext in SUPPORTED_EXTENSIONS:
        for filepath in glob.glob(os.path.join(docs_path, "**", f"*{ext}"), recursive=True):
            try:
                rel_parts = Path(filepath).resolve().relative_to(docs_p).parts
                if any(p.startswith(".") or p in ("__pycache__", "node_modules", "venv", "env", "chroma_db", "build", "dist") for p in rel_parts):
                    continue
                all_files.append(filepath)
            except ValueError:
                pass # Path not relative to docs_path
    all_files = sorted(set(all_files))  # deduplicate

    if not all_files:
        return f"No supported files found in {docs_path}\nSupported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"

    # ── Incremental indexing: detect changed files ──
    hash_cache = _load_hash_cache() if not force_reindex else {}
    changed_files = []
    unchanged_count = 0
    new_hash_cache = {}

    for filepath in all_files:
        try:
            current_hash = _get_file_hash(filepath)
        except Exception:
            continue
        norm_path = os.path.normpath(filepath)
        new_hash_cache[norm_path] = current_hash
        if not force_reindex and hash_cache.get(norm_path) == current_hash:
            unchanged_count += 1
        else:
            changed_files.append(filepath)

    # Detect deleted files (in cache but no longer on disk)
    current_paths = {os.path.normpath(f) for f in all_files}
    deleted_paths = set(hash_cache.keys()) - current_paths

    # If nothing changed and nothing deleted, skip
    if not changed_files and not deleted_paths:
        return (
            f"Nothing changed. {unchanged_count} files already up to date.\n"
            f"Collection: {collection}\n"
            f"Path: {docs_path}\n"
            f"Use force_reindex=True to rebuild everything."
        )

    target_col = _get_collection(collection)

    # ── Handle deleted files: remove their chunks ──
    if deleted_paths:
        try:
            for del_path in deleted_paths:
                del_filename = Path(del_path).name
                old_data = target_col.get(
                    where={"filename": del_filename},
                    include=[]
                )
                if old_data["ids"]:
                    target_col.delete(ids=old_data["ids"])
        except Exception:
            pass

    # ── If force or first-time: full re-index ──
    if force_reindex or target_col.count() == 0:
        try:
            client.delete_collection(collection)
        except Exception:
            pass
        target_col = _get_collection(collection)
        changed_files = all_files

    # ── Remove old chunks from changed files ──
    if not force_reindex:
        for filepath in changed_files:
            fname = Path(filepath).name
            try:
                old_data = target_col.get(
                    where={"filename": fname},
                    include=[]
                )
                if old_data["ids"]:
                    target_col.delete(ids=old_data["ids"])
            except Exception:
                pass

    # ── Index changed files ──
    ids, texts, metas = [], [], []

    for filepath in changed_files:
        raw = _read_file_to_text(filepath)
        if not raw:
            continue

        category = _categorize_file(filepath, raw)
        sections = _extract_sections_smart(raw, filepath)

        for sec in sections:
            path_hash = hashlib.md5(filepath.encode("utf-8")).hexdigest()[:6]
            chunk_hash = hashlib.md5(sec["text"].encode("utf-8")).hexdigest()[:10]
            chunk_id = f"{Path(filepath).stem}_{path_hash}__{chunk_hash}"
            chunk_id = re.sub(r"[^a-zA-Z0-9_]", "_", chunk_id)

            ids.append(chunk_id)
            texts.append(sec["text"])
            metas.append({
                "source": sec["source"],
                "filename": sec["filename"],
                "heading": sec["heading"],
                "parent_heading": sec["parent_heading"],
                "category": category,
                "word_count": sec["word_count"]
            })

    # ── Add to ChromaDB in batches ──
    batch_size = 50
    added = 0
    for i in range(0, len(ids), batch_size):
        batch_end = min(i + batch_size, len(ids))
        target_col.add(
            ids=ids[i:batch_end],
            documents=texts[i:batch_end],
            metadatas=metas[i:batch_end],
        )
        added += batch_end - i

    # ── Rebuild BM25 from full collection ──
    global _bm25_loaded_from
    all_data = target_col.get(include=["documents"])
    all_ids = all_data["ids"]
    all_texts = all_data["documents"]
    _build_bm25(all_ids, all_texts)
    _bm25_loaded_from = "indexed"

    # Save hash cache
    _save_hash_cache(new_hash_cache)

    categories = {}
    for m in metas:
        cat = m["category"]
        categories[cat] = categories.get(cat, 0) + 1

    cat_summary = " | ".join(f"{k}: {v}" for k, v in sorted(categories.items()))
    total_chunks = target_col.count()

    # Count file types
    ext_counts = {}
    for fp in changed_files:
        ext = Path(fp).suffix.lower()
        ext_counts[ext] = ext_counts.get(ext, 0) + 1
    ext_str = ", ".join(f"{e}: {c}" for e, c in sorted(ext_counts.items()))

    status_parts = []
    if deleted_paths:
        status_parts.append(f"Deleted: {len(deleted_paths)} removed files")
    status_parts.append(
        f"Indexed {added} chunks from {len(changed_files)} changed files "
        f"({unchanged_count} unchanged)"
    )

    return (
        f"{' | '.join(status_parts)}\n"
        f"Total: {total_chunks} chunks in collection '{collection}'\n"
        f"File types: {ext_str}\n"
        f"Path: {docs_path}\n"
        f"Categories: {cat_summary}\n"
        f"Model: {EMBED_MODEL} | Device: {DEVICE}\n"
        f"Reranker: {RERANK_MODEL}\n"
        f"DB: {DB_PATH}"
    )


def _expand_query(query: str) -> list[str]:
    """
    Generate query variations to improve recall.
    Especially important for RU documentation with EN API terms.
    """
    queries = [query]
    has_cyrillic = any('\u0400' <= c <= '\u04FF' for c in query)

    if has_cyrillic:
        queries.append(query.lower())
        api_synonyms = {
            "создать": "create post make new",
            "получить": "get fetch list read",
            "обновить": "update patch put",
            "удалить": "delete remove drop",
            "найти": "find search query",
            "добавить": "add insert append",
            "ошибка": "error exception fail bug",
        }
        en_terms = []
        for ru, en in api_synonyms.items():
            if ru in query.lower():
                en_terms.append(en)
        if en_terms:
            queries.append(" ".join(en_terms))

    return queries[:3]


def _deduplicate_results(candidates: dict) -> dict:
    """Remove near-duplicates (>80% similar by first 200 chars)."""
    from difflib import SequenceMatcher

    items = sorted(candidates.values(), key=lambda x: x.get("final_score", 0), reverse=True)
    kept = []
    for item in items:
        item_text = item["doc"][:200]
        is_dup = any(
            SequenceMatcher(None, item_text, k["doc"][:200]).ratio() > 0.8
            for k in kept
        )
        if not is_dup:
            kept.append(item)
    return {str(i): item for i, item in enumerate(kept)}


@mcp.tool()
def search_docs(
    query: str,
    n_results: int = 5,
    category: Optional[str] = None,
    filename: Optional[str] = None,
    collection: str = DEFAULT_COLLECTION
) -> str:
    """
    Semantic search across indexed Markdown documentation.
    Returns the most relevant chunks with source file references.

    Args:
        query: Natural language query, e.g. "how does the authentication logic work?"
        n_results: Number of results to return (default 5)
        category: Optional filter by category (dynamically based on your root folder names)
        filename: Optional filter by specific file, e.g. "architecture.md"
        collection: Collection to search in (default: docs_v4)
    """
    target_col = _get_collection(collection)
    count = target_col.count()
    if count == 0:
        return f"No documents indexed in collection '{collection}'. Call index_documents() first."

    # Validate inputs
    n_results = max(1, n_results)

    # ── Step 1: Vector search (ChromaDB) with Query Expansion ──
    where_filter = None
    conditions = []
    if category:
        conditions.append({"category": category})
    if filename:
        conditions.append({"filename": filename})

    if len(conditions) == 1:
        where_filter = conditions[0]
    elif len(conditions) > 1:
        where_filter = {"$and": conditions}

    fetch_count = min(n_results * 4, count)  # fetch more for reranking
    expanded_queries = _expand_query(query)
    candidates: dict[str, dict] = {}

    try:
        for q in expanded_queries:
            vector_results = target_col.query(
                query_texts=[q],
                n_results=fetch_count,
                where=where_filter,
                include=["documents", "metadatas", "distances"]
            )
            
            if not vector_results["documents"][0]:
                continue
                
            result_ids = vector_results.get("ids", [[]])[0]
            for i, (doc, meta, dist) in enumerate(zip(
                vector_results["documents"][0],
                vector_results["metadatas"][0],
                vector_results["distances"][0]
            )):
                chunk_id = result_ids[i] if i < len(result_ids) else f"_vec_{i}"
                if chunk_id not in candidates:
                    candidates[chunk_id] = {
                        "doc": doc, "meta": meta,
                        "vec_score": 1 - dist, "bm25_score": 0.0,
                        "min_vec_rank": i
                    }
                else:
                    candidates[chunk_id]["min_vec_rank"] = min(
                        candidates[chunk_id]["min_vec_rank"], i
                    )
    except Exception as e:
        return f"Search error: {str(e)}"

    if not candidates:
        return f"No results found for: '{query}'"

    # ── Step 2: BM25 keyword search ──
    if _bm25_index is not None and _bm25_corpus is not None:
        query_tokens = _tokenize(query)
        if query_tokens:
            bm25_scores = _bm25_index.get_scores(query_tokens)
            max_bm25 = max(bm25_scores) if max(bm25_scores) > 0 else 1.0

            # Match BM25 results to vector candidates by chunk_id
            for idx, (chunk_id, text) in enumerate(_bm25_corpus):
                if chunk_id in candidates:
                    candidates[chunk_id]["bm25_score"] = bm25_scores[idx] / max_bm25

    # ── Step 3: Combine scores (RRF - Reciprocal Rank Fusion) ──
    # Sort by vector rank (min_vec_rank from expansion)
    vec_sorted = sorted(candidates.values(), key=lambda x: x.get("min_vec_rank", 999))
    for rank, c in enumerate(vec_sorted):
        c["vec_rank"] = rank + 1

    # Sort by BM25 score to get BM25 rank
    bm25_sorted = sorted(candidates.values(), key=lambda x: x["bm25_score"], reverse=True)
    for rank, c in enumerate(bm25_sorted):
        c["bm25_rank"] = rank + 1

    # RRF: combined_score = 1/(k+vec_rank) + 1/(k+bm25_rank)
    k = 60  # standard RRF constant
    for c in candidates.values():
        c["hybrid_score"] = (1.0 / (k + c["vec_rank"])) + (1.0 / (k + c["bm25_rank"]))

    # Sort by hybrid score
    hybrid_sorted = sorted(candidates.values(), key=lambda x: x["hybrid_score"], reverse=True)

    # Take top candidates for reranking
    top_candidates = hybrid_sorted[:min(n_results * 2, len(hybrid_sorted))]

    # ── Step 4: Cross-Encoder Reranking ──
    try:
        cross_encoder = _get_cross_encoder()
        pairs = [(query, c["doc"][:MAX_CROSS_ENCODER_CHARS]) for c in top_candidates]
        ce_scores = cross_encoder.predict(pairs)

        # Normalize CE scores to 0-1
        min_ce = min(ce_scores)
        max_ce = max(ce_scores) if max(ce_scores) != min(ce_scores) else min(ce_scores) + 1
        for i, c in enumerate(top_candidates):
            c["ce_score"] = (ce_scores[i] - min_ce) / (max_ce - min_ce)

        # Final score: 40% hybrid + 60% cross-encoder
        # Pre-compute min/max outside the loop (was O(n²), now O(n))
        h_scores = [c["hybrid_score"] for c in top_candidates]
        max_hybrid, min_hybrid = max(h_scores), min(h_scores)
        hybrid_range = max_hybrid - min_hybrid if max_hybrid != min_hybrid else 1.0

        for c in top_candidates:
            norm_hybrid = (c["hybrid_score"] - min_hybrid) / hybrid_range
            c["final_score"] = 0.4 * norm_hybrid + 0.6 * c["ce_score"]

    except Exception:
        # Fallback: use hybrid score directly
        if top_candidates:
            max_h = max(c["hybrid_score"] for c in top_candidates)
            for c in top_candidates:
                c["final_score"] = c["hybrid_score"] / max_h if max_h > 0 else 0

    # Sort by final score and apply deduplication
    final_sorted = sorted(top_candidates, key=lambda x: x["final_score"], reverse=True)
    deduped_candidates = _deduplicate_results({str(i): c for i, c in enumerate(final_sorted)})
    final_results = list(deduped_candidates.values())[:n_results]

    # ── Step 5: Format output ──
    output_parts = [f"## Results for: \"{query}\"\n"]

    if category or filename:
        filters = []
        if category:
            filters.append(f"category={category}")
        if filename:
            filters.append(f"file={filename}")
        output_parts.append(f"**Filters:** {', '.join(filters)}\n")

    output_parts.append(f"**Method:** Hybrid (Vector + BM25) + Cross-Encoder Reranking\n")

    for rank, c in enumerate(final_results, 1):
        output_parts.append(_format_result(c["doc"], c["meta"], c["final_score"], rank))

    return "\n".join(output_parts)


@mcp.tool()
def rag_status(collection: str = DEFAULT_COLLECTION) -> str:
    """Show how many chunks are indexed and from which files.

    Args:
        collection: Collection to inspect (default: docs_v4)
    """
    target_col = _get_collection(collection)
    count = target_col.count()
    if count == 0:
        return f"No documents indexed in collection '{collection}'. Run index_documents() to start."

    all_data = target_col.get(include=["metadatas"])["metadatas"]

    files = {}
    categories = {}
    total_words = 0
    for m in all_data:
        fname = m.get("filename", "unknown")
        cat = m.get("category", "other")
        wc = m.get("word_count", 0)
        files[fname] = files.get(fname, 0) + 1
        categories[cat] = categories.get(cat, 0) + 1
        total_words += wc

    file_list = "\n".join(
        f"  - {name} ({chunks} chunks)"
        for name, chunks in sorted(files.items())
    )
    cat_list = "\n".join(
        f"  - {cat}: {n} chunks"
        for cat, n in sorted(categories.items())
    )

    if _bm25_index is not None:
        bm25_status = f"Active ({len(_bm25_corpus)} chunks, from {_bm25_loaded_from or 'unknown'})"
    else:
        bm25_status = "Not built (re-index needed)"
    ce_status = "Loaded" if _cross_encoder is not None else "Lazy (loads on first search)"

    return (
        f"## RAG Index Status\n\n"
        f"**Collection:** {collection}\n"
        f"**Total:** {count} chunks | ~{total_words:,} words\n"
        f"**Embedding:** {EMBED_MODEL} ({DEVICE})\n"
        f"**BM25:** {bm25_status}\n"
        f"**Cross-Encoder:** {ce_status} ({RERANK_MODEL})\n"
        f"**DB:** {DB_PATH}\n\n"
        f"### Files ({len(files)}):\n{file_list}\n\n"
        f"### Categories:\n{cat_list}"
    )


@mcp.tool()
def list_collections() -> str:
    """List all RAG collections in the database with their chunk counts."""
    try:
        collections = client.list_collections()
    except Exception as e:
        return f"Error listing collections: {e}"

    if not collections:
        return "No collections found. Run index_documents() to create one."

    lines = ["## RAG Collections\n"]
    for col in collections:
        try:
            count = col.count()
            lines.append(f"- **{col.name}** — {count} chunks")
        except Exception:
            lines.append(f"- **{col.name}** — (error reading)")
    return "\n".join(lines)


@mcp.tool()
def list_indexed_files(collection: str = DEFAULT_COLLECTION) -> str:
    """List all files that have been indexed in a collection.

    Args:
        collection: Collection to inspect (default: docs_v4)
    """
    target_col = _get_collection(collection)
    count = target_col.count()
    if count == 0:
        return f"No documents in collection '{collection}'."

    all_meta = target_col.get(include=["metadatas"])["metadatas"]

    files = {}
    for m in all_meta:
        fname = m.get("filename", "unknown")
        source = m.get("source", "")
        if fname not in files:
            files[fname] = {"chunks": 0, "source": source, "category": m.get("category", "")}
        files[fname]["chunks"] += 1

    lines = [f"## Indexed Files in '{collection}' ({len(files)} files, {count} chunks)\n"]
    for fname, info in sorted(files.items()):
        lines.append(f"- **{fname}** — {info['chunks']} chunks | category: {info['category']}")
        if info["source"]:
            lines.append(f"  Source: `{info['source']}`")
    return "\n".join(lines)


@mcp.tool()
def remove_source(
    filename: str,
    collection: str = DEFAULT_COLLECTION
) -> str:
    """Remove all chunks for a specific file from the index.

    Args:
        filename: Filename to remove, e.g. "api.md"
        collection: Collection to remove from (default: docs_v4)
    """
    target_col = _get_collection(collection)

    try:
        data = target_col.get(
            where={"filename": filename},
            include=[]
        )
    except Exception as e:
        return f"Error: {e}"

    if not data["ids"]:
        return f"No chunks found for file '{filename}' in collection '{collection}'."

    count = len(data["ids"])
    target_col.delete(ids=data["ids"])

    # Rebuild BM25 after removal
    global _bm25_loaded_from
    remaining = target_col.get(include=["documents"])
    _build_bm25(remaining["ids"], remaining["documents"])
    _bm25_loaded_from = "indexed"

    return f"Removed {count} chunks for '{filename}' from collection '{collection}'."


@mcp.tool()
def delete_collection(
    name: str,
    confirm: bool = False
) -> str:
    """Delete an entire collection. Requires confirm=True for safety.

    Args:
        name: Collection name to delete
        confirm: Must be True to actually delete. Without it, shows what would be deleted.
    """
    try:
        target_col = _get_collection(name)
        chunk_count = target_col.count()
    except Exception:
        chunk_count = 0

    if not confirm:
        return (
            f"⚠️ Will delete collection '{name}' ({chunk_count} chunks).\n"
            f"To confirm: delete_collection(name='{name}', confirm=True)"
        )

    try:
        client.delete_collection(name)
    except Exception as e:
        return f"Error deleting collection '{name}': {e}"

    # Rebuild BM25 (collection is gone, may need cleanup)
    global _bm25_loaded_from
    try:
        default_col = _get_collection(DEFAULT_COLLECTION)
        remaining = default_col.get(include=["documents"])
        _build_bm25(remaining["ids"], remaining["documents"])
        _bm25_loaded_from = "indexed"
    except Exception:
        pass

    return f"✅ Deleted collection '{name}' ({chunk_count} chunks removed)."


@mcp.tool()
def reindex_collection(
    docs_path: str = DOCS_PATH,
    collection: str = DEFAULT_COLLECTION
) -> str:
    """Force full reindex of a collection. Deletes all existing chunks and re-indexes from scratch.

    Args:
        docs_path: Path to scan for files (default: parent directory)
        collection: Collection name to reindex (default: docs_v4)
    """
    return index_documents(docs_path=docs_path, collection=collection, force_reindex=True)


# ──────────────────────────────────────────────
# Web Crawling Integration (Phase 4)
# ──────────────────────────────────────────────

def _index_web_pages(
    pages: list[dict],
    collection: str,
    source_label: str,
) -> str:
    """
    Bridge function: takes crawled pages and indexes them into ChromaDB.
    Reuses the existing chunking pipeline from server.py.
    """
    target_col = _get_collection(collection)

    ids, texts, metas = [], [], []

    for page in pages:
        content = page["content"]
        url = page.get("url", "web")

        # Use the existing chunking pipeline
        sections = _extract_sections(content, url)
        category = _categorize_file(url, content)

        for sec in sections:
            url_hash = hashlib.md5(url.encode("utf-8")).hexdigest()[:6]
            chunk_hash = hashlib.md5(sec["text"].encode("utf-8")).hexdigest()[:10]
            chunk_id = f"web_{url_hash}__{chunk_hash}"
            chunk_id = re.sub(r"[^a-zA-Z0-9_]", "_", chunk_id)

            ids.append(chunk_id)
            texts.append(sec["text"])
            metas.append({
                "source": url,
                "filename": url.split("/")[-1] or "index",
                "heading": sec["heading"],
                "parent_heading": sec["parent_heading"],
                "category": category,
                "word_count": sec["word_count"],
            })

    if not ids:
        return f"No content extracted from {source_label}."

    # Add to ChromaDB in batches
    batch_size = 50
    added = 0
    for i in range(0, len(ids), batch_size):
        batch_end = min(i + batch_size, len(ids))
        target_col.add(
            ids=ids[i:batch_end],
            documents=texts[i:batch_end],
            metadatas=metas[i:batch_end],
        )
        added += batch_end - i

    # Rebuild BM25
    global _bm25_loaded_from
    all_data = target_col.get(include=["documents"])
    _build_bm25(all_data["ids"], all_data["documents"])
    _bm25_loaded_from = "indexed"

    return (
        f"Indexed {added} chunks from {len(pages)} pages\n"
        f"Source: {source_label}\n"
        f"Collection: {collection}\n"
        f"Total: {target_col.count()} chunks"
    )


@mcp.tool()
def index_url(
    uri: str,
    collection: str = DEFAULT_COLLECTION,
    max_pages: int = 200,
    max_depth: int = 10,
    stay_within_prefix: bool = True,
    exclude_patterns: list[str] = None,
    use_sitemap: bool = True,
    use_playwright: bool = False,
) -> str:
    """
    Index any web source into the RAG database. Supports:

    - URLs:    index_url("https://docs.python.org/3/library/")
    - GitHub:  index_url("github://owner/repo") or index_url("github://owner/repo/docs")
    - npm:     index_url("npm://axios@1.6")
    - PyPI:    index_url("pypi://fastapi")
    - ZIP:     index_url("file:///C:/path/to/docs.zip")

    Args:
        uri:                  Source URI to index
        collection:           Collection to index into (default: docs_v4)
        max_pages:            Max pages to crawl for URLs (default: 200)
        max_depth:            Max link depth from starting URL (default: 10)
        stay_within_prefix:   Don't leave the starting URL path (default: True)
        exclude_patterns:     Regex patterns to skip URLs (e.g. ["/blog/", "/changelog/"])
        use_sitemap:          Try sitemap.xml first for faster discovery (default: True)
        use_playwright:       Use headless browser for JS-rendered sites (default: False)
    """
    import asyncio

    try:
        from crawler import (
            crawl_and_index,
            index_github,
            index_npm,
            index_pypi,
            index_zip,
        )
    except ImportError as e:
        return (
            f"Error: crawler module not available. {e}\n"
            "Make sure crawler.py is in the same directory as server.py\n"
            "and install: pip install httpx beautifulsoup4 html2text"
        )

    try:
        try:
            asyncio.get_running_loop()
            is_running = True
        except RuntimeError:
            is_running = False

        if is_running:
            import concurrent.futures
            def _run_in_thread():
                return asyncio.run(_run_index(
                    uri, collection, max_pages, max_depth,
                    stay_within_prefix, exclude_patterns,
                    use_sitemap, use_playwright,
                ))
            with concurrent.futures.ThreadPoolExecutor() as pool:
                future = pool.submit(_run_in_thread)
                return future.result()
        else:
            return asyncio.run(_run_index(
                uri, collection, max_pages, max_depth,
                stay_within_prefix, exclude_patterns,
                use_sitemap, use_playwright,
            ))
    except Exception as e:
        return f"Error during indexing: {e}"


async def _run_index(
    uri: str,
    collection: str,
    max_pages: int,
    max_depth: int,
    stay_within_prefix: bool,
    exclude_patterns: list[str] | None,
    use_sitemap: bool,
    use_playwright: bool,
) -> str:
    """Async dispatcher for index_url()."""
    from crawler import (
        crawl_and_index,
        index_github,
        index_npm,
        index_pypi,
        index_zip,
    )

    if uri.startswith(("http://", "https://")):
        pages, status = await crawl_and_index(
            uri, collection, max_pages, max_depth,
            stay_within_prefix, exclude_patterns,
            use_sitemap, use_playwright,
        )
        label = f"Web: {uri}"

    elif uri.startswith("github://"):
        pages, status = await index_github(uri, collection)
        label = f"GitHub: {uri}"

    elif uri.startswith("npm://"):
        pages, status = await index_npm(uri, collection)
        label = f"npm: {uri}"

    elif uri.startswith("pypi://"):
        pages, status = await index_pypi(uri, collection)
        label = f"PyPI: {uri}"

    elif uri.lower().endswith(".zip"):
        local_path = uri.replace("file:///", "").replace("file://", "")
        pages, status = await index_zip(local_path)
        label = f"ZIP: {uri}"

    else:
        return f"Error: unknown URI scheme '{uri}'. Use http://, https://, github://, npm://, pypi://"

    if not pages:
        return f"No content found.\n{status}"

    index_result = _index_web_pages(pages, collection, label)
    return f"{status}\n{index_result}"


if __name__ == "__main__":
    # Start web dashboard if configured
    if os.getenv("RAG_DASHBOARD", "").lower() in ("true", "1", "yes"):
        try:
            from dashboard import start_dashboard_thread
            start_dashboard_thread()
        except Exception:
            pass  # dashboard is optional

    # Start file watcher if configured (FEAT-09)
    _watcher_observer = None
    if os.getenv("RAG_WATCH_PATH"):
        try:
            from watcher import start_watcher
            _watcher_observer = start_watcher(
                os.getenv("RAG_WATCH_PATH"),
                os.getenv("RAG_WATCH_COLLECTION", DEFAULT_COLLECTION),
            )
        except Exception:
            pass  # watcher is optional
    mcp.run(transport="stdio")
